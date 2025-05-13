import streamlit as st
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from PIL import Image
import io
import os
from pathlib import Path
import base64
import numpy as np
import cohere
import google.generativeai as genai
from dotenv import load\_dotenv
import tempfile
import shutil
from langchain.schema import HumanMessage, AIMessage

# Load environment variables

load\_dotenv()

# Initialize APIs

cohere\_api\_key = os.getenv("COHERE\_API\_KEY")
gemini\_api\_key = os.getenv("GOOGLE\_API\_KEY")

if not cohere\_api\_key or not gemini\_api\_key:
st.error("API keys not found. Please check your .env file")
st.stop()

try:
co = cohere.Client(api\_key=cohere\_api\_key)
genai.configure(api\_key=gemini\_api\_key)
except Exception as e:
st.error(f"Failed to initialize API clients: {str(e)}")
st.stop()

# Configuration

MAX\_PIXELS = 1568 \* 1568
SUPPORTED\_TYPES = \["pdf", "docx", "pptx"]
GEMINI\_MODEL = "gemini-1.5-flash"  # Current recommended model

# Create temporary directory

OUTPUT\_DIR = Path(tempfile.mkdtemp())
IMAGES\_DIR = OUTPUT\_DIR / "images"
TEXT\_FILE = OUTPUT\_DIR / "extracted\_text.txt"

def cleanup():
"""Remove temporary files"""
if OUTPUT\_DIR.exists():
shutil.rmtree(OUTPUT\_DIR)

def save\_image(image\_pil, image\_count):
"""Save image to temporary directory"""
IMAGES\_DIR.mkdir(parents=True, exist\_ok=True)
img\_path = IMAGES\_DIR / f"image\_{image\_count}.png"
image\_pil.save(img\_path)
return str(img\_path)

def resize\_image(pil\_image):
"""Resize image if too large"""
org\_width, org\_height = pil\_image.size
if org\_width \* org\_height > MAX\_PIXELS:
scale\_factor = (MAX\_PIXELS / (org\_width \* org\_height)) \*\* 0.5
new\_width = int(org\_width \* scale\_factor)
new\_height = int(org\_height \* scale\_factor)
pil\_image.thumbnail((new\_width, new\_height))

def base64\_from\_image(img\_path):
"""Convert image to base64"""
try:
pil\_image = Image.open(img\_path)
img\_format = pil\_image.format or "PNG"
resize\_image(pil\_image)
with io.BytesIO() as buffer:
pil\_image.save(buffer, format=img\_format)
encoded = base64.b64encode(buffer.getvalue()).decode("utf-8")
return f"data\:image/{img\_format.lower()};base64,{encoded}"
except Exception as e:
st.error(f"Error processing image: {str(e)}")
return None

def extract\_pdf(file):
"""Extract text and images from PDF"""
text = ""
image\_paths = \[]
image\_count = 1

```
try:
    with fitz.open(stream=file.read(), filetype="pdf") as pdf:
        for page in pdf:
            text += page.get_text()
            for img in page.get_images(full=True):
                xref = img[0]
                base_image = pdf.extract_image(xref)
                img_bytes = base_image["image"]
                img_pil = Image.open(io.BytesIO(img_bytes))
                img_path = save_image(img_pil, image_count)
                image_paths.append(img_path)
                image_count += 1
except Exception as e:
    st.error(f"Error processing PDF: {str(e)}")
return text, image_paths
```

def extract\_docx(file):
"""Extract text and images from DOCX"""
text = ""
image\_paths = \[]
image\_count = 1

```
try:
    doc = Document(file)
    for para in doc.paragraphs:
        text += para.text + "\n"

    for rel in doc.part._rels:
        rel_obj = doc.part._rels[rel]
        if "image" in rel_obj.target_ref:
            img_data = rel_obj.target_part.blob
            img_pil = Image.open(io.BytesIO(img_data))
            img_path = save_image(img_pil, image_count)
            image_paths.append(img_path)
            image_count += 1
except Exception as e:
    st.error(f"Error processing DOCX: {str(e)}")
return text, image_paths
```

def extract\_pptx(file):
"""Extract text and images from PPTX"""
text = ""
image\_paths = \[]
image\_count = 1

```
try:
    prs = Presentation(file)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
            if shape.shape_type == 13:  # Picture shape
                img_stream = shape.image.blob
                img_pil = Image.open(io.BytesIO(img_stream))
                img_path = save_image(img_pil, image_count)
                image_paths.append(img_path)
                image_count += 1
except Exception as e:
    st.error(f"Error processing PPTX: {str(e)}")
return text, image_paths
```

def ask\_gemini(question, context=None, img\_path=None):
"""Query Gemini with optional context and/or image"""
try:
model = genai.GenerativeModel(GEMINI\_MODEL)

```
    if img_path and context:
        # Both image and text context
        prompt = f"""Answer the question based on the following image and additional context if relevant.
```

If the question is general and not related to the context, answer it generally.

Context: {context}
Question: {question}"""
img = Image.open(img\_path)
response = model.generate\_content(\[prompt, img])
elif img\_path:
\# Image only
prompt = f"""Answer the question based on the following image if relevant.
If the question is general and not about the image, answer it generally.

Question: {question}"""
img = Image.open(img\_path)
response = model.generate\_content(\[prompt, img])
elif context:
\# Text context only
prompt = f"""Answer the question based on the following context if relevant.
If the question is general and not related to the context, answer it generally.

Context: {context}
Question: {question}"""
response = model.generate\_content(prompt)
else:
\# No context - general question
response = model.generate\_content(question)

```
    return response.text
except Exception as e:
    return f"Error querying Gemini: {str(e)}"
```

# Streamlit UI

st.set\_page\_config(page\_title="Document Q\&A", layout="wide")
st.title("📄 Document Q\&A with Image Analysis")

# Initialize session state for chat histories

if 'text\_chat\_history' not in st.session\_state:
st.session\_state.text\_chat\_history = \[]
if 'image\_chat\_history' not in st.session\_state:
st.session\_state.image\_chat\_history = \[]
if 'general\_chat\_history' not in st.session\_state:
st.session\_state.general\_chat\_history = \[]
if 'scroll' not in st.session\_state:
st.session\_state.scroll = False
if 'processed' not in st.session\_state:
st.session\_state.processed = False
st.session\_state.text = ""
st.session\_state.image\_paths = \[]
st.session\_state.selected\_img = None

# File upload section

with st.expander("Upload Document", expanded=True):
uploaded\_file = st.file\_uploader(
"Choose a file",
type=SUPPORTED\_TYPES,
help=f"Supported formats: {', '.join(SUPPORTED\_TYPES)}"
)

# Process file when uploaded

if uploaded\_file and not st.session\_state.processed:
with st.spinner("Extracting content from document..."):
file\_ext = uploaded\_file.name.split(".")\[-1].lower()

```
    try:
        if file_ext == "pdf":
            st.session_state.text, st.session_state.image_paths = extract_pdf(uploaded_file)
        elif file_ext == "docx":
            st.session_state.text, st.session_state.image_paths = extract_docx(uploaded_file)
        elif file_ext == "pptx":
            st.session_state.text, st.session_state.image_paths = extract_pptx(uploaded_file)
        
        # Save text to file
        OUTPUT_DIR.mkdir(exist_ok=True)
        with open(TEXT_FILE, "w", encoding="utf-8") as f:
            f.write(st.session_state.text)
        
        st.session_state.processed = True
        st.success("Document processed successfully!")
    except Exception as e:
        st.error(f"Failed to process document: {str(e)}")
        cleanup()
```

# Create tabs for different chat interfaces

tab1, tab2, tab3 = st.tabs(\["📝 Text Analysis", "🖼️ Image Analysis", "💬 General Chat"])

def render\_chat(container, chat\_history):
"""Render chat messages in a container"""
with container:
for message in chat\_history:
if isinstance(message, HumanMessage):
st.markdown(
f"<div style='text-align: right; color: white; background-color: #0a84ff; padding: 8px; border-radius: 10px; margin: 5px 0; max-width: 80%; float: right; clear: both;'>{message.content}</div>",
unsafe\_allow\_html=True
)
elif isinstance(message, AIMessage):
st.markdown(
f"<div style='text-align: left; color: black; background-color: #d1d1d1; padding: 8px; border-radius: 10px; margin: 5px 0; max-width: 80%; float: left; clear: both;'>{message.content}</div>",
unsafe\_allow\_html=True
)

# Text Analysis Tab

with tab1:
st.subheader("Text Analysis")

```
# Display extracted text
if st.session_state.processed:
    with st.expander("View Extracted Text"):
        st.text_area("Extracted Text", st.session_state.text, height=200, label_visibility="collapsed")

# Chat interface for text analysis
text_chat_container = st.container()

user_text_input = st.text_input(
    "Ask about the text content:", 
    key="text_input",
    placeholder="Type your question here...",
    label_visibility="collapsed"
)
text_send_button = st.button("Send", key="text_send")

if text_send_button and user_text_input:
    st.session_state.text_chat_history.append(HumanMessage(content=user_text_input))
    if user_text_input.lower() == 'close the chat':
        st.stop()
    
    with st.spinner("Analyzing text..."):
        answer = ask_gemini(user_text_input, context=st.session_state.text)
        st.session_state.text_chat_history.append(AIMessage(content=answer))
        st.session_state.scroll = True
        st.rerun()

render_chat(text_chat_container, st.session_state.text_chat_history)
```

# Image Analysis Tab

# Image Analysis Tab

with tab2:
st.subheader("Image Analysis")

```
if st.session_state.processed and st.session_state.image_paths:
    st.write("Select an image to analyze:")
    
    selected_img_index = None
    num_cols = 3
    image_paths = st.session_state.image_paths
    rows = (len(image_paths) + num_cols - 1) // num_cols

    for row in range(rows):
        cols = st.columns(num_cols)
        for col_idx in range(num_cols):
            img_idx = row * num_cols + col_idx
            if img_idx < len(image_paths):
                img_path = image_paths[img_idx]
                with cols[col_idx]:
                    try:
                        img = Image.open(img_path)
                        img.thumbnail((200, 200))
                        st.image(img, use_container_width=True)
                        if st.button(f"Analyze Image {img_idx+1}", key=f"btn_{img_idx}"):
                            st.session_state.selected_img = img_path
                    except Exception as e:
                        st.error(f"Error loading image: {str(e)}")

    if st.session_state.selected_img:
        st.divider()
        st.subheader("Analyzing Selected Image")

        selected_img = Image.open(st.session_state.selected_img)
        selected_img.thumbnail((300, 300))  # Resize selected image to be smaller
        st.image(selected_img, caption="Selected Image")

        image_chat_container = st.container()
        user_image_input = st.text_input(
            "Ask about the image:", 
            key="image_input",
            placeholder="Type your question here...",
            label_visibility="collapsed"
        )
        image_send_button = st.button("Send", key="image_send")

        if image_send_button and user_image_input:
            st.session_state.image_chat_history.append(HumanMessage(content=user_image_input))
            if user_image_input.lower() == 'close the chat':
                st.stop()
            with st.spinner("Analyzing image..."):
                answer = ask_gemini(user_image_input, img_path=st.session_state.selected_img, context=st.session_state.text)
                st.session_state.image_chat_history.append(AIMessage(content=answer))
                st.session_state.scroll = True
                st.rerun()

        render_chat(image_chat_container, st.session_state.image_chat_history)
    else:
        st.info("Please select an image to begin analysis.")
else:
    st.write("No images found in the document.")
```

# General Chat Tab

with tab3:
st.subheader("General Chat")

```
# Chat interface for general questions
general_chat_container = st.container()

user_general_input = st.text_input(
    "Ask any question:", 
    key="general_input",
    placeholder="Type your question here...",
    label_visibility="collapsed"
)
general_send_button = st.button("Send", key="general_send")

if general_send_button and user_general_input:
    st.session_state.general_chat_history.append(HumanMessage(content=user_general_input))
    if user_general_input.lower() == 'close the chat':
        st.stop()
    
    with st.spinner("Thinking..."):
        answer = ask_gemini(user_general_input)
        st.session_state.general_chat_history.append(AIMessage(content=answer))
        st.session_state.scroll = True
        st.rerun()

render_chat(general_chat_container, st.session_state.general_chat_history)
```

# Auto-scroll to bottom of chat

if st.session\_state.scroll:
st.session\_state.scroll = False
st.markdown(
""" <script>
window\.scrollTo(0, document.body.scrollHeight); </script>
""",
unsafe\_allow\_html=True
)

# Clean up when done

st.session\_state.cleanup = cleanup   
